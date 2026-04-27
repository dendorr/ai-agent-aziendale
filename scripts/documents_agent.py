"""
DOCUMENTS AGENT v3.1 (ASYNC) — Universal document intelligence
Converts ALL documents to structured Markdown before indexing.
Optimized for multi-user architecture and vLLM/OpenAI server.
"""

import sys, os, json, re, hashlib, asyncio
from pathlib import Path

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from config.config import (CHROMA_PATHS, LLM_MODEL, LLM_BASE_URL, LLM_API_KEY,
                           CHUNK_SIZE, CHUNK_OVERLAP, EXTENSIONS, MEMORY_PATH)
import chromadb
import semantic_analyzer as analyzer
from openai import AsyncOpenAI

client     = chromadb.PersistentClient(path=CHROMA_PATHS["documents"])
collection = client.get_or_create_collection("documents")

MEMORY_FILE    = Path(MEMORY_PATH) / "documents_memory.json"
MARKDOWN_CACHE = Path(MEMORY_PATH).parent / "markdown_cache"

# ── Models ─────────────────────────────────────────────────────────────────────
ROUTING_MODEL = "qwen3:0.6b"  
ANSWER_MODEL  = LLM_MODEL      

# Initialize shared async OpenAI client
llm_client = AsyncOpenAI(base_url=LLM_BASE_URL, api_key=LLM_API_KEY)

# ── OCR — singleton to avoid re-initializing easyocr every call ───────────────
_easyocr_reader = None

def _get_easyocr_reader():
    """Lazy-initialize the easyocr Reader once per process."""
    global _easyocr_reader
    if _easyocr_reader is None:
        try:
            import easyocr
            _easyocr_reader = easyocr.Reader(
                ["it", "en"], gpu=False, verbose=False
            )
        except Exception:
            pass
    return _easyocr_reader

def ocr_image_bytes(image_bytes: bytes, source_hint: str = "") -> str:
    """Run OCR on raw image bytes. Level 1: pytesseract. Level 2: easyocr."""
    if not image_bytes:
        return ""
    from PIL import Image
    import io
    try:
        img = Image.open(io.BytesIO(image_bytes))
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
    except Exception:
        return ""

    # Level 1: pytesseract
    try:
        import pytesseract
        text = pytesseract.image_to_string(
            img, lang="ita+eng", config="--psm 3 --oem 3"
        ).strip()
        if len(text) > 15:
            return f"[OCR]\n{text}"
    except Exception:
        pass

    # Level 2: easyocr
    try:
        import numpy as np
        reader = _get_easyocr_reader()
        if reader:
            results = reader.readtext(np.array(img))
            lines   = [r[1] for r in results if len(r) > 2 and r[2] > 0.3]
            text    = "\n".join(lines).strip()
            if len(text) > 15:
                return f"[OCR]\n{text}"
    except Exception:
        pass

    return source_hint

def ocr_image_file(filepath) -> str:
    """OCR on a file path."""
    try:
        return ocr_image_bytes(Path(filepath).read_bytes())
    except Exception:
        return ""

# ── Shared helpers ─────────────────────────────────────────────────────────────

def _table_to_markdown(table) -> str:
    """Convert python-pptx or python-docx Table to Markdown table."""
    rows = []
    for i, row in enumerate(table.rows):
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
        if i == 0: 
            rows.append("|" + "|".join(["---"] * len(cells)) + "|")
    return "\n".join(rows)

# ── PPTX → Markdown ────────────────────────────────────────────────────────────

def convert_pptx_to_markdown(filepath) -> str:
    """Convert a PPTX file to structured Markdown."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    p   = Path(filepath)
    prs = Presentation(str(filepath))

    lines = [f"# {p.stem}", f"*File: {p.name} — {len(prs.slides)} slides*", ""]

    for slide_num, slide in enumerate(prs.slides, 1):
        title = ""
        try:
            if slide.shapes.title and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()
        except Exception:
            pass

        slide_heading = f"## Slide {slide_num}"
        if title:
            slide_heading += f": {title}"
        lines.append(slide_heading)
        lines.append("")

        for shape in slide.shapes:
            if shape.has_table:
                lines.append(_table_to_markdown(shape.table))
                lines.append("")
                continue

            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text or text == title:
                        continue
                    level  = getattr(para, "level", 0)
                    indent = "  " * level
                    prefix = f"{'#' * (level + 3)} " if level > 0 else ""
                    lines.append(f"{indent}{prefix}{text}")
                lines.append("")
                continue

            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_bytes = shape.image.blob
                    hint      = f"[Image — slide {slide_num}]"
                    ocr_text  = ocr_image_bytes(img_bytes, source_hint=hint)
                    if ocr_text:
                        quoted = ocr_text.replace("\n", "\n> ")
                        lines.append(f"> **{hint}**")
                        lines.append(f"> {quoted}")
                        lines.append("")
            except Exception:
                pass

        try:
            notes_tf   = slide.notes_slide.notes_text_frame
            notes_text = notes_tf.text.strip()
            if notes_text:
                quoted = notes_text.replace("\n", "\n> ")
                lines.append("**Speaker Notes:**")
                lines.append(f"> {quoted}")
                lines.append("")
        except Exception:
            pass

        lines.append("---")
        lines.append("")

    return "\n".join(lines)

# ── DOCX → Markdown ────────────────────────────────────────────────────────────

def convert_docx_to_markdown(filepath) -> str:
    """Convert a DOCX file to structured Markdown."""
    p = Path(filepath)

    try:
        from markitdown import MarkItDown
        result = MarkItDown().convert(str(filepath))
        text   = result.text_content.strip()
        if len(text) > 100:
            md = f"# {p.stem}\n*File: {p.name}*\n\n{text}"
            img_blocks = _docx_extract_image_ocr(filepath)
            if img_blocks:
                md += "\n\n## Extracted Images\n\n" + "\n\n".join(img_blocks)
            return md
    except ImportError:
        pass
    except Exception as e:
        print(f"  [markitdown warning] {p.name}: {e}", flush=True)

    try:
        from docx import Document
        from docx.oxml.ns import qn
        from docx.text.paragraph import Paragraph
        from docx.table import Table as DocxTable

        doc   = Document(str(filepath))
        lines = [f"# {p.stem}", f"*File: {p.name}*", ""]

        _HEADING_MAP = {
            "heading 1": "#", "heading 2": "##", "heading 3": "###",
            "heading 4": "####", "heading 5": "#####",
        }

        for child in doc.element.body.iterchildren():
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

            if tag == "p":
                para  = Paragraph(child, doc)
                text  = para.text.strip()
                style = (para.style.name or "").lower() if para.style else ""

                if not text:
                    lines.append("")
                    continue

                prefix = ""
                for key, md_prefix in _HEADING_MAP.items():
                    if key in style:
                        prefix = md_prefix
                        break

                if prefix:
                    lines.append(f"{prefix} {text}")
                elif "list" in style:
                    level  = getattr(para, "style", None)
                    indent = "  " * (int(re.search(r'\d', style).group()) - 1
                                     if re.search(r'\d', style) else 0)
                    lines.append(f"{indent}- {text}")
                else:
                    lines.append(text)

            elif tag == "tbl":
                tbl = DocxTable(child, doc)
                lines.append("")
                lines.append(_table_to_markdown(tbl))
                lines.append("")

        lines.append("")
        img_blocks = _docx_extract_image_ocr(filepath)
        if img_blocks:
            lines.append("## Extracted Images")
            lines.extend(img_blocks)

        return "\n".join(lines)

    except Exception as e:
        return f"[DOCX error: {e}]"

def _docx_extract_image_ocr(filepath) -> list:
    """Extract embedded images from DOCX and run OCR."""
    blocks = []
    try:
        from docx import Document
        doc = Document(str(filepath))
        for i, rel in enumerate(doc.part.rels.values(), 1):
            if "image" in rel.reltype:
                try:
                    img_bytes = rel.target_part.blob
                    ocr_text  = ocr_image_bytes(
                        img_bytes, source_hint=f"[Image {i}]"
                    )
                    if ocr_text and len(ocr_text) > 10:
                        quoted = ocr_text.replace("\n", "\n> ")
                        blocks.append(f"> **Image {i}:**\n> {quoted}")
                except Exception:
                    pass
    except Exception:
        pass
    return blocks

# ── PDF → Markdown ─────────────────────────────────────────────────────────────

def convert_pdf_to_markdown(filepath) -> str:
    """Convert a PDF to structured Markdown."""
    import fitz as _fitz
    p      = Path(filepath)
    header = f"# {p.stem}\n*File: {p.name}*\n\n"

    try:
        from markitdown import MarkItDown
        result = MarkItDown().convert(str(filepath))
        text   = result.text_content.strip()
        if len(text) > 100:
            return header + text
    except ImportError:
        pass
    except Exception as e:
        print(f"  [markitdown warning] {p.name}: {e}", flush=True)

    fitz_parts = []
    try:
        doc = _fitz.open(str(filepath))
        for page_num, page in enumerate(doc, 1):
            page_lines = [f"## Page {page_num}", ""]

            text = page.get_text("text")
            if text.strip():
                page_lines.append(text.strip())
                page_lines.append("")

            try:
                for tab in page.find_tables().tables:
                    import pandas as pd
                    df = tab.to_pandas()
                    md = df.to_markdown(index=False)
                    if md:
                        page_lines.append(f"**Table:**\n{md}")
                        page_lines.append("")
            except Exception:
                pass

            for img_info in page.get_images(full=True):
                try:
                    xref      = img_info[0]
                    base_img  = doc.extract_image(xref)
                    img_bytes = base_img.get("image", b"")
                    ocr_text  = ocr_image_bytes(
                        img_bytes, source_hint=f"[Image page {page_num}]"
                    )
                    if ocr_text and len(ocr_text) > 15:
                        quoted = ocr_text.replace("\n", "\n> ")
                        page_lines.append(f"> **Image page {page_num}:**\n> {quoted}")
                        page_lines.append("")
                except Exception:
                    pass

            fitz_parts.append("\n".join(page_lines))
        doc.close()
    except Exception as e:
        fitz_parts.append(f"[fitz error: {e}]")

    plumber_parts = []
    try:
        import pdfplumber
        with pdfplumber.open(str(filepath)) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                for table in page.extract_tables() or []:
                    rows = []
                    for i, row in enumerate(table):
                        cells = [str(c).strip() if c else "" for c in row]
                        rows.append("| " + " | ".join(cells) + " |")
                        if i == 0:
                            rows.append("|" + "|".join(["---"] * len(cells)) + "|")
                    if rows:
                        plumber_parts.append(
                            f"**Table page {page_num}:**\n" + "\n".join(rows)
                        )
    except Exception:
        pass

    result = header
    if fitz_parts:
        result += "\n\n---\n\n".join(fitz_parts)
    if plumber_parts:
        result += "\n\n## Additional Tables\n\n" + "\n\n".join(plumber_parts)

    return result if len(result) > len(header) + 20 else f"[Extraction failed for {p.name}]"

# ── Markdown cache ─────────────────────────────────────────────────────────────

def _cache_path(filepath) -> Path:
    h    = hashlib.md5(str(filepath).encode()).hexdigest()[:8]
    stem = re.sub(r"[^a-zA-Z0-9]", "_", Path(filepath).stem[:40]).strip("_")
    return MARKDOWN_CACHE / f"{stem}_{h}.md"

def _cache_is_valid(filepath, cache_file: Path) -> bool:
    if not cache_file.exists():
        return False
    try:
        return cache_file.stat().st_mtime >= Path(filepath).stat().st_mtime
    except Exception:
        return False

def get_or_create_markdown(filepath) -> str:
    """Reads from disk cache if valid; otherwise converts and writes the cache."""
    MARKDOWN_CACHE.mkdir(parents=True, exist_ok=True)
    cache_file = _cache_path(filepath)

    if _cache_is_valid(filepath, cache_file):
        try:
            return cache_file.read_text(encoding="utf-8")
        except Exception:
            pass

    ext      = Path(filepath).suffix.lower()
    markdown = ""

    if ext == ".pdf":
        markdown = convert_pdf_to_markdown(filepath)
    elif ext in (".pptx", ".ppt"):
        markdown = convert_pptx_to_markdown(filepath)
    elif ext in (".docx", ".doc"):
        markdown = convert_docx_to_markdown(filepath)
    elif ext == ".md":
        markdown = Path(filepath).read_text(encoding="utf-8", errors="ignore")
    elif ext == ".txt":
        content  = Path(filepath).read_text(encoding="utf-8", errors="ignore")
        markdown = f"# {Path(filepath).stem}\n\n{content}"

    if markdown and len(markdown.strip()) > 50:
        try:
            cache_file.write_text(markdown, encoding="utf-8")
            print(f"  [cache] Written {cache_file.name} "
                  f"({len(markdown):,} chars)", flush=True)
        except Exception as e:
            print(f"  [cache warning] {e}", flush=True)

    return markdown

# ── Memory ─────────────────────────────────────────────────────────────────────

def load_memory() -> dict:
    if MEMORY_FILE.exists():
        try:
            return json.loads(MEMORY_FILE.read_text(encoding="utf-8"))
        except Exception:
            pass
    return {"documents": {}, "annotations": {}}

def save_memory(memory: dict):
    MEMORY_FILE.parent.mkdir(parents=True, exist_ok=True)
    MEMORY_FILE.write_text(
        json.dumps(memory, indent=2, ensure_ascii=False), encoding="utf-8"
    )

def update_document_memory(filepath, markdown: str):
    from datetime import datetime
    memory = load_memory()
    if "documents" not in memory:
        memory["documents"] = {}

    p = Path(filepath)
    memory["documents"][p.name] = {
        "filepath":   str(filepath),
        "indexed_at": datetime.now().isoformat(timespec="seconds"),
        "words":      len(markdown.split()),
        "chars":      len(markdown),
        "size_kb":    round(p.stat().st_size / 1024, 1),
        "ext":        p.suffix.lower(),
    }
    save_memory(memory)

def add_annotation(filename: str, note: str):
    memory = load_memory()
    if "annotations" not in memory:
        memory["annotations"] = {}
    if filename not in memory["annotations"]:
        memory["annotations"][filename] = []
    memory["annotations"][filename].append(note)
    save_memory(memory)

# ── Chunking + indexing ────────────────────────────────────────────────────────

def chunk_text(text: str) -> list:
    words = text.split()
    chunks, i = [], 0
    while i < len(words):
        chunks.append(" ".join(words[i:i + CHUNK_SIZE]))
        i += CHUNK_SIZE - CHUNK_OVERLAP
    return chunks or [""]

def index_file(filepath) -> int:
    p   = Path(filepath)
    ext = p.suffix.lower()

    if ext not in EXTENSIONS["documents"]:
        return 0

    print(f"  [docs] Processing {p.name}...", flush=True)
    markdown = get_or_create_markdown(filepath)

    if not markdown or not markdown.strip():
        return 0

    update_document_memory(filepath, markdown)

    try:
        existing = collection.get(where={"path": str(filepath)})
        if existing and existing["ids"]:
            collection.delete(ids=existing["ids"])
    except Exception:
        pass

    chunks = chunk_text(markdown)
    for i, chunk in enumerate(chunks):
        collection.upsert(
            documents=[chunk],
            ids=[f"{filepath}__c{i}"],
            metadatas=[{
                "filename": p.name,
                "path":     str(filepath),
                "chunk":    i,
                "agent":    "documents",
                "type":     "chunk",
                "ext":      ext,
            }]
        )
    return len(chunks)

def index_folder(folder):
    files = [f for f in Path(folder).rglob("*")
             if f.is_file() and f.suffix.lower() in EXTENSIONS["documents"]]
    print(f"[Documents] Found {len(files)} files...")
    total = 0
    for f in files:
        n = index_file(str(f))
        if n: print(f"  [OK] {f.name} → {n} chunks")
        else: print(f"  [--] {f.name} → skipped")
        total += n
    print(f"[Documents] Completed — {total} chunks total")

# ── Search ─────────────────────────────────────────────────────────────────────

def search(query: str):
    """Synchronous vector DB search (handled via threads in server.py)."""
    return analyzer.search_with_cards(collection, query, "documents", n_results=6)

# ── Routing model (ASYNC) ──────────────────────────────────────────────────────

async def route_documents(query: str, memory: dict) -> list:
    """Async routing to identify relevant documents without blocking."""
    docs = memory.get("documents", {})
    if not docs:
        return []

    doc_lines = "\n".join(
        f"  - {name}  ({info.get('words', 0):,} words | "
        f"{info.get('size_kb', 0)} KB | {info.get('ext', '')})"
        for name, info in list(docs.items())[:25]
    )

    prompt = f"""Sei un assistente che seleziona i documenti più rilevanti per rispondere a una domanda.
DOCUMENTI DISPONIBILI:
{doc_lines}
DOMANDA: {query}
Elenca SOLO i nomi file più rilevanti (massimo 4), uno per riga. Solo nomi file esatti, nessun altro testo."""

    try:
        response = await llm_client.chat.completions.create(
            model=ROUTING_MODEL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.0
        )
        raw = response.choices[0].message.content
        all_docs = set(docs.keys())
        relevant = [
            line.strip().strip("-").strip()
            for line in raw.strip().split("\n")
            if line.strip() and any(d in line for d in all_docs)
        ]
        return relevant[:4]
    except Exception as e:
        print(f"Routing doc error: {e}")
        return []

# ── System prompt (Italian as requested by domain) ─────────────────────────────

SYSTEM_PROMPT = """Sei un esperto analista documentale con accesso completo ai documenti aziendali.

REGOLE ASSOLUTE:
- Rispondi SEMPRE e SOLO in italiano, qualunque sia la lingua del documento
- Usa ESCLUSIVAMENTE le informazioni presenti nei documenti forniti nel contesto
- Non inventare, non stimare, non integrare con conoscenze esterne
- Cita sempre la fonte (nome file, numero slide, numero pagina) per ogni dato riportato
- Se l'informazione non è presente nei documenti, dichiaralo esplicitamente

NUMERI, MISURE E SPECIFICHE TECNICHE — REGOLA CRITICA:
- Riporta i valori ESATTAMENTE come appaiono nel documento: nessuna trasformazione
- Preserva le unità di misura (mm, cm, m, µm, kg, g, ml, l, %, °C, bar, N, pz, €, $, ecc.)
- Non arrotondare mai, non convertire unità, non cambiare il formato
- Tolleranze (±0,5; ±5%; max 3 mm) vanno sempre riportate insieme al valore principale
- Tabelle con misure vanno riportate COMPLETE e FEDELI all'originale
- Se il documento riporta un range (es. 10–15 kg), riporta il range intero

TESTO ESTRATTO DA IMMAGINI ([OCR]):
- Il testo tra [OCR] proviene da immagini, grafici o schemi nel documento
- Trattalo come dato attendibile; segnala la fonte OCR se utile al contesto

STRUTTURA DELLE RISPOSTE:
- Specifiche tecniche → tabella o lista ordinata con unità
- Confronti tra documenti → colonne affiancate con fonte per ogni valore
- Riassunti → struttura gerarchica (sezioni → punti chiave)
- Ricerca di valori specifici → cita il contesto esatto del documento

CAPACITÀ:
- Analisi di presentazioni, relazioni, manuali, capitolati, specifiche tecniche
- Estrazione precisa di dati da tabelle, grafici e immagini (via OCR)
- Confronto tra più documenti sullo stesso argomento
- Ricerca di termini, misure, codici o specifiche esatte
- Sintesi di documenti lunghi mantenendo tutti i dati numerici"""

# ── Answer (ASYNC) ─────────────────────────────────────────────────────────────

async def answer(question: str, context: str) -> str:
    """Async answer generation optimized for concurrency."""
    memory   = load_memory()
    docs     = memory.get("documents", {})
    mem_txt  = ""

    if docs:
        doc_lines = "\n".join(
            f"  - {name} ({info.get('words', 0):,} parole | {info.get('ext', '')})"
            for name, info in list(docs.items())[:20]
        )
        mem_txt += f"\nDOCUMENTI INDICIZZATI:\n{doc_lines}\n"

    annotations = memory.get("annotations", {})
    if annotations:
        ann_lines = [
            f"  [{fname}] {note}"
            for fname, notes in list(annotations.items())[:5]
            for note in notes[:3]
        ]
        if ann_lines:
            mem_txt += "\nANNOTAZIONI UTENTE:\n" + "\n".join(ann_lines) + "\n"

    # Await the async routing call
    relevant = await route_documents(question, memory)
    if relevant:
        mem_txt += "\nDOCUMENTI PROBABILMENTE RILEVANTI PER QUESTA DOMANDA:\n"
        mem_txt += "\n".join(f"  - {d}" for d in relevant) + "\n"

    prompt = f"""{SYSTEM_PROMPT}
{mem_txt}
=== CONTENUTO DOCUMENTI — FONTE PRIMARIA ===
{context}

DOMANDA: {question}

RISPOSTA (obbligatoriamente in italiano — dati solo dalla FONTE PRIMARIA,
numeri e misure ESATTAMENTE come nel documento):"""

    try:
        # Generate final response async via OpenAI compatible API
        response = await llm_client.chat.completions.create(
            model=ANSWER_MODEL,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.1
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"LLM Connection Error: {e}"

if __name__ == "__main__":
    folder = sys.argv[1] if len(sys.argv) > 1 else "."
    index_folder(folder)